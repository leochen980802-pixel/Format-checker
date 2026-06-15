import streamlit as st
import google.generativeai as genai
import docx
import fitz  # PyMuPDF
from pptx import Presentation
import json

import json
import os
from datetime import datetime
import streamlit as st
# (確保你原本的 import genai, time 等等都在)

# 定義儲存歷史紀錄的檔案名稱
HISTORY_FILE = "proofread_history.json"

# --- 讀取歷史紀錄 ---
def load_history():
    if os.path.exists(HISTORY_FILE):
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                try:
                                        return json.load(f)
                except json.JSONDecodeError:
                    return [] # 如果檔案損壞，回傳空陣列
                    return []

# --- 寫入歷史紀錄 ---
                    def save_history(file_name, results):
                        history = load_history()
                        new_record = {
                        "id": datetime.now().strftime("%Y%m%d%H%M%S"), # 產生唯一 ID
                        "time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "file_name": file_name,
                        "results": results
                        }
# 把最新的紀錄加到陣列最前面
history.insert(0, new_record) 
                                                                                                                                
with open(HISTORY_FILE, "w", encoding="utf-8") as f:
                                                                                                                                            json.dump(history, f, ensure_ascii=False, indent=4)


# --- 1. 設定與初始化 ---
st.set_page_config(page_title="學術格式校對 AI", layout="centered")

# 從 Streamlit 雲端環境安全地取得 API Key
try:
    api_key = st.secrets["GEMINI_API_KEY"]
    genai.configure(api_key=api_key)
except Exception as e:
    st.error("系統找不到 API 金鑰。請確認你已經在 Streamlit 的 Advanced settings 中設定了 GEMINI_API_KEY。")
    st.stop()

# --- 2. 核心功能：讀取不同格式的文件 ---
def extract_text(file, file_type):
    text = ""
    try:
        if file_type == "docx":
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == "pdf":
            doc = fitz.open(stream=file.read(), filetype="pdf")
            for page in doc:
                text += page.get_text()
        elif file_type == "pptx":
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        st.error(f"檔案讀取失敗：{e}")
    return text

# --- 3. 核心功能：呼叫 AI 進行校對 ---
import time  # 務必確保檔案最上方有 import time

# --- 3. 核心功能：呼叫 AI 進行校對（商用規格版） ---
def proofread_text(text, format_style):
    # 預設使用 Flash 模型，以取得每分鐘 15 次的高呼叫額度
    target_model = "gemini-1.5-flash" 
    
    try:
        for m in genai.list_models():
            if 'generateContent' in m.supported_generation_methods:
                if 'flash' in m.name:
                    target_model = m.name
                    break
    except Exception:
        pass 

# --- 側邊欄：歷史紀錄區塊 ---
st.sidebar.title("🗂️ 校對歷史紀錄")
history_data = load_history()

# 使用 Session State 來追蹤使用者目前想看哪一筆歷史紀錄
if "selected_history" not in st.session_state:
    st.session_state.selected_history = None

if not history_data:
    st.sidebar.info("目前還沒有任何校對紀錄。")
else:
    for record in history_data:
        # 在側邊欄為每一筆紀錄產生一個按鈕
        button_label = f"📄 {record['file_name']} \n({record['time']})"
        if st.sidebar.button(button_label, key=record['id']):
            # 當按下按鈕時，把該筆紀錄存入 session_state
            st.session_state.selected_history = record['results']
            st.session_state.current_viewing_file = record['file_name']

# 如果使用者有點擊歷史紀錄，就在主畫面顯示出來，並提供一個「返回」按鈕
if st.session_state.selected_history:
    st.info(f"回顧模式：正在查看【{st.session_state.current_viewing_file}】的校對結果")
    if st.button("⬅️ 返回新增校對"):
        st.session_state.selected_history = None
        st.rerun() # 重新載入畫面
        
    # 將歷史紀錄的結果渲染出來 (這裡請套用你原本顯示錯誤清單的邏輯)
    for item in st.session_state.selected_history:
        st.error(f"❌ 原文：{item.get('original', '')}")
        st.warning(f"💡 問題：{item.get('issue', '')}")
        st.success(f"✅ 建議：{item.get('fix', '')}")
        st.markdown("---")
        
    st.stop() # 停止執行下方的上傳介面，保持在回顧模式


    try:
            model = genai.GenerativeModel(target_model)
    except Exception as e:
            st.error(f"模型載入失敗：{e}")

    # =================【核心升級 1：滑動視窗分段處理】=================
    # 將 64 頁的龐大文字，每 1500 字切成一塊，確保 AI 能集中注意力「逐字」細看
    chunk_size = 1500
    text_chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    
    all_results = []
    
    # 在 Streamlit 畫面上建立進度條與狀態文字
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    # 開始分批輪詢校對
    for idx, chunk in enumerate(text_chunks):
        status_text.text(f"正在深度校對第 {idx+1} / {len(text_chunks)} 段簡報文字...")
        
        prompt = f"""
        你現在是一位擁有 20 年經驗的嚴苛學術期刊主編。你的任務是進行極度精準的文字與格式校對。
        請嚴格執行以下步驟：
        1. 逐字掃描，找出「所有」錯別字、漏字與不通順的文法。
        2. 嚴格檢查格式是否「完全」符合【{format_style}】規範（包含引註格式、標點符號全半形、大小寫等）。
        3. 寧可嚴格，不可錯漏。必須抓出所有問題。

        你必須嚴格回傳一個 JSON 陣列，格式如下：
        [
          {{"original": "錯誤的句子或單字", "issue": "具體的錯誤原因說明", "fix": "建議的精確修改內容"}}
        ]
        
        待校對文字段落：
        {chunk} 
        """
        
        try:
            # =================【核心升級 2：原生 JSON 模式】=================
            # 透過 response_mime_type 強制 API 只能輸出純 JSON，徹底根除 Extra data 錯誤
            response = model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.0,
                    "response_mime_type": "application/json"
                }
            )
            
            # 直接解析，不需要再寫 find('[') 暴力切字串
            chunk_result = json.loads(response.text)
            if isinstance(chunk_result, list):
                all_results.extend(chunk_result)
                
        except Exception as e:
            error_msg = str(e)
            # =================【核心升級 3：自動對抗 429 限制】=================
            if "429" in error_msg or "Quota" in error_msg:
                status_text.text("⚠️ 觸發免費版頻率限制，系統自動等待 30 秒後繼續...")
                time.sleep(30)
                # 重試機制
                try:
                    response = model.generate_content(
                        prompt, 
                        generation_config={"temperature": 0.0, "response_mime_type": "application/json"}
                    )
                    chunk_result = json.loads(response.text)
                    if isinstance(chunk_result, list):
                        all_results.extend(chunk_result)
                except Exception:
                    pass
            else:
                st.error(f"第 {idx+1} 段解析失敗，已跳過。錯誤：{error_msg}")
        
        # =================【核心升級 4：調頻緩衝】=================
        # 免費版 Flash 限制每分鐘 15 次，每次呼叫完刻意休息 4 秒，確保整體運作平穩不中斷
        time.sleep(4)
        progress_bar.progress((idx + 1) / len(text_chunks))
        
    status_text.text("✨ 全本 64 頁簡報深度校對完成！")
    time.sleep(1)
    progress_bar.empty()
    status_text.empty()
    

#def proofread_text(text, format_style):
    # 【策略調整】退回 Flash 模型以獲得每分鐘 15 次的扣打，適合頻繁除錯測試
    target_model = "gemini-1.5-flash" 
    
    try:
        # 動態尋找 Flash 模型
        for m in genai.list_models():
            if 'generateContent' in m.supported_generation_methods:
                if 'flash' in m.name:
                    target_model = m.name
                    break
    except Exception:
        pass 

    try:
        model = genai.GenerativeModel(target_model)
    except Exception as e:
        st.error(f"模型載入失敗：{e}")

    # 提示詞保持不變（維持專家身分與嚴格度）
    prompt = f"""
    你現在是一位擁有 20 年經驗的嚴苛學術期刊主編。你的任務是進行極度精準的文字與格式校對。
    請嚴格執行以下步驟：
    1. 逐字掃描，找出「所有」錯別字、漏字與不通順的文法。
    2. 嚴格檢查格式是否「完全」符合【{format_style}】規範（包含引註格式、標點符號全半形、大小寫等）。
    3. 寧可嚴格，不可錯漏。必須抓出所有問題。

    請務必以 JSON 陣列格式回傳，絕對不要包含任何其他文字、問候語或 Markdown 標記，格式嚴格如下：
    [
      {{"original": "錯誤的句子或單字", "issue": "具體的錯誤原因說明", "fix": "建議的精確修改內容"}}
    ]
    
    待校對文字：
    {text[:4000]} 
    """
    
    try:
        # 溫度同樣維持 0.0，讓 Flash 也能像機器一樣精準執行
        response = model.generate_content(
            prompt,
            generation_config={"temperature": 0.0}
        )
        result_text = response.text
        
        # 暴力萃取法
        start_idx = result_text.find('[')
        end_idx = result_text.rfind(']')
        
        if start_idx != -1 and end_idx != -1:
            clean_json = result_text[start_idx:end_idx+1]
        else:
            st.error("AI 未回傳標準的 JSON 格式，請再試一次。")
            
    except Exception as e:
        # 【優化錯誤捕捉】如果是 429，給予白話文提示而不是滿江紅
        error_msg = str(e)
        if "429" in error_msg or "Quota" in error_msg:
             st.warning("⚠️ 呼叫太頻繁啦！免費版 API 有頻率限制，請等待約 1 分鐘後再點擊測試。")
        else:
             st.error(f"AI 解析資料失敗，錯誤訊息：{error_msg}")


# --- 4. 介面與互動邏輯 ---
st.title("📄 學術格式與錯字校對工具")
st.markdown("上傳你的 Word, PDF 或 PPT，AI 將自動抓出格式瑕疵與錯字。")

# --- 側邊欄：歷史紀錄區塊 ---
st.sidebar.title("🗂️ 校對歷史紀錄")
history_data = load_history()

# 使用 Session State 來追蹤使用者目前想看哪一筆歷史紀錄
if "selected_history" not in st.session_state:
    st.session_state.selected_history = None

if not history_data:
    st.sidebar.info("目前還沒有任何校對紀錄。")
else:
    for record in history_data:
        # 在側邊欄為每一筆紀錄產生一個按鈕
        button_label = f"📄 {record['file_name']} \n({record['time']})"
        if st.sidebar.button(button_label, key=record['id']):
            # 當按下按鈕時，把該筆紀錄存入 session_state
            st.session_state.selected_history = record['results']
            st.session_state.current_viewing_file = record['file_name']

# 如果使用者有點擊歷史紀錄，就在主畫面顯示出來，並提供一個「返回」按鈕
if st.session_state.selected_history:
    st.info(f"回顧模式：正在查看【{st.session_state.current_viewing_file}】的校對結果")
    if st.button("⬅️ 返回新增校對"):
        st.session_state.selected_history = None
        st.rerun() # 重新載入畫面
        
    # 將歷史紀錄的結果渲染出來 (這裡請套用你原本顯示錯誤清單的邏輯)
    for item in st.session_state.selected_history:
        st.error(f"❌ 原文：{item.get('original', '')}")
        st.warning(f"💡 問題：{item.get('issue', '')}")
        st.success(f"✅ 建議：{item.get('fix', '')}")
        st.markdown("---")
        
    st.stop() # 停止執行下方的上傳介面，保持在回顧模式


st.sidebar.header("設定")
format_choice = st.sidebar.selectbox(
    "請選擇目標格式規範：",
    ["APA 格式", "MLA 格式", "Chicago 格式"]
)

uploaded_file = st.file_uploader(
    "上傳文件 (支援 .docx, .pdf, .pptx)", 
    type=["docx", "pdf", "pptx"]
)

if uploaded_file is not None:
    # 取得檔案副檔名
    file_ext = uploaded_file.name.split(".")[-1].lower()
    st.success(f"已成功上傳：{uploaded_file.name}")
    
    if st.button("開始 AI 校對"):
        with st.spinner("AI 正在努力閱讀與校對中，這可能需要幾十秒，請稍候..."):
            
            # 1. 將上傳的檔案轉換為純文字
            raw_text = extract_text(uploaded_file, file_ext)
            
            if not raw_text.strip():
                st.warning("無法從檔案中讀取到文字，請確認檔案內容並非純圖片。")
            else:
            # 2. 將文字與選擇的格式送給 API
                issues = proofread_text(raw_text, format_choice)
                
            # 3. 將陣列中的錯誤一條一條列印在畫面上
            if issues is None:
                # 如果是 None，代表上面已經印出紅字錯誤了，這裡什麼都不做 (pass)
                pass
            elif len(issues) > 0:
                    st.subheader("📝 發現以下格式或錯字問題：")
                    for idx, item in enumerate(issues):
                        with st.expander(f"問題 {idx + 1}：{item.get('issue', '格式問題')}"):
                            st.write("**原文：**", item.get("original", ""))
                            st.write("**建議修改：**", item.get("fix", ""))
            else:
                    # 只有真正回傳了空陣列 []，才代表完全沒錯字
                    st.success("太棒了！AI 沒有發現明顯的格式錯誤或錯字。")

